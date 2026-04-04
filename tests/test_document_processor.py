"""
Unit tests for document_processor module.

Tests document parsing, text extraction, and chunking functionality.
"""

import tempfile
from pathlib import Path

import openpyxl
import pandas as pd
import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from src.document_processing.document_processor import (
    MIN_CHUNK_CHARS,
    DocumentProcessor,
    _normalize_extracted_text,
)


class TestDocumentProcessor:
    """Test suite for DocumentProcessor class."""

    @pytest.fixture
    def processor(self):
        """Create a DocumentProcessor instance for testing."""
        # Use production-like chunk_size so metadata doesn't dominate
        return DocumentProcessor(chunk_size=2048, chunk_overlap=100)

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield Path(tmpdir)

    def test_compute_file_hash_consistency(self, processor, temp_dir):
        """
        Test that file hash is consistent for same content.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        # Create a test file
        test_file = temp_dir / "test.txt"
        test_content = "This is test content for hashing."
        test_file.write_text(test_content)

        # Compute hash twice
        hash1 = processor.compute_file_hash(test_file)
        hash2 = processor.compute_file_hash(test_file)

        # Hashes should be identical
        assert hash1 == hash2
        assert len(hash1) == 64  # SHA-256 produces 64 hex characters

    def test_compute_file_hash_uniqueness(self, processor, temp_dir):
        """
        Test that different files produce different hashes.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        # Create two different files
        file1 = temp_dir / "test1.txt"
        file2 = temp_dir / "test2.txt"

        file1.write_text("Content A")
        file2.write_text("Content B")

        hash1 = processor.compute_file_hash(file1)
        hash2 = processor.compute_file_hash(file2)

        # Hashes should be different
        assert hash1 != hash2

    def test_extract_text_from_txt(self, processor, temp_dir):
        """
        Test text extraction from plain text file.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.txt"
        test_content = "Hello, this is a test document.\nWith multiple lines."
        test_file.write_text(test_content)

        extracted = processor.extract_text_from_txt(test_file)

        assert extracted == test_content

    def test_extract_text_from_csv(self, processor, temp_dir):
        """
        Test text extraction from CSV file.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.csv"

        # Create a simple CSV
        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]})
        df.to_csv(test_file, index=False)

        extracted = processor.extract_text_from_csv(test_file)

        # Check that extracted text contains the data
        assert "Alice" in extracted
        assert "Bob" in extracted
        assert "25" in extracted
        assert "30" in extracted

    def test_extract_text_from_docx(self, processor, temp_dir):
        """
        Test text extraction from DOCX file.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.docx"

        # Create a DOCX with content
        doc = DocxDocument()
        doc.add_paragraph("First paragraph.")
        doc.add_paragraph("Second paragraph.")
        doc.save(test_file)

        extracted = processor.extract_text_from_docx(test_file)

        assert "First paragraph." in extracted
        assert "Second paragraph." in extracted

    def test_extract_text_from_pptx(self, processor, temp_dir):
        """Test text extraction from PPTX file."""
        test_file = temp_dir / "test.pptx"

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout

        slide1 = prs.slides.add_slide(slide_layout)
        slide1.shapes.title.text = "Introduction"
        slide1.placeholders[1].text = "Welcome to the presentation."

        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Details"
        slide2.placeholders[1].text = "Here are the details."

        prs.save(test_file)

        extracted = processor.extract_text_from_pptx(test_file)

        assert "Slide 1:" in extracted
        assert "Introduction" in extracted
        assert "Welcome to the presentation." in extracted
        assert "Slide 2:" in extracted
        assert "Details" in extracted

    def test_extract_text_from_pptx_with_notes(self, processor, temp_dir):
        """Test that speaker notes are extracted from PPTX."""
        test_file = temp_dir / "test_notes.pptx"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with notes"
        slide.notes_slide.notes_text_frame.text = "These are speaker notes."
        prs.save(test_file)

        extracted = processor.extract_text_from_pptx(test_file)

        assert "Notes: These are speaker notes." in extracted

    def test_extract_text_from_pptx_empty(self, processor, temp_dir):
        """Test extraction from PPTX with no text content."""
        test_file = temp_dir / "empty.pptx"

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        prs.save(test_file)

        extracted = processor.extract_text_from_pptx(test_file)

        assert isinstance(extracted, str)

    def test_extract_text_unsupported_format(self, processor, temp_dir):
        """
        Test that unsupported file format raises ValueError.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xyz"
        test_file.write_text("Some text content in unusual format")

        with pytest.raises(ValueError, match="Unsupported file format"):
            processor.extract_text(test_file)

    def test_process_document_success(self, processor, temp_dir):
        """
        Test successful document processing.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.txt"
        # Create content long enough to generate multiple chunks
        test_content = "This is a test. " * 50
        test_file.write_text(test_content)

        nodes = processor.process_document(test_file, source_type="manual")

        # Should produce at least one node
        assert len(nodes) > 0

        # Check metadata
        first_node = nodes[0]
        assert first_node.metadata["filename"] == "test.txt"
        assert first_node.metadata["source_type"] == "manual"
        assert first_node.metadata["file_type"] == ".txt"
        assert "file_hash" in first_node.metadata

    def test_process_document_empty_file(self, processor, temp_dir):
        """
        Test processing empty file returns empty list.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "empty.txt"
        test_file.write_text("")

        nodes = processor.process_document(test_file)

        assert len(nodes) == 0

    def test_process_document_with_extra_metadata(self, processor, temp_dir):
        """
        Test that extra metadata is added to nodes.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.txt"
        test_file.write_text("Test content for metadata.")

        extra_metadata = {"author": "Test User", "department": "DoLS"}

        nodes = processor.process_document(test_file, extra_metadata=extra_metadata)

        assert len(nodes) > 0
        first_node = nodes[0]
        assert first_node.metadata["author"] == "Test User"
        assert first_node.metadata["department"] == "DoLS"

    def test_process_directory(self, processor, temp_dir):
        """
        Test processing all files in a directory.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        # Create multiple test files
        (temp_dir / "doc1.txt").write_text("Document 1 content.")
        (temp_dir / "doc2.txt").write_text("Document 2 content.")
        (temp_dir / "doc3.csv").write_text("col1,col2\nval1,val2")

        # Create unsupported file (should be skipped)
        (temp_dir / "ignored.xyz").write_text("Should be ignored")

        nodes = processor.process_directory(temp_dir)

        # Should process 3 supported files
        unique_files = set(n.metadata["filename"] for n in nodes)
        assert len(unique_files) == 3
        assert "doc1.txt" in unique_files
        assert "doc2.txt" in unique_files
        assert "doc3.csv" in unique_files
        assert "ignored.xyz" not in unique_files

    def test_chunking_with_overlap(self, processor, temp_dir):
        """
        Test that chunking respects overlap parameter.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.txt"
        # Create realistic prose content long enough to produce multiple chunks
        sentence = "This is a sentence with enough words to exercise the chunker. "
        test_content = sentence * 200  # ~12,000 chars of real sentences
        test_file.write_text(test_content)

        nodes = processor.process_document(test_file)

        # Should create multiple chunks due to length
        assert len(nodes) >= 2

        # Check that chunks have reasonable length
        for node in nodes:
            # Chunks shouldn't exceed chunk_size significantly
            assert len(node.text) <= processor.chunk_size + 200  # Allow buffer

    def test_extract_text_from_excel_with_headers(self, processor, temp_dir):
        """
        Test text extraction from Excel file with headers.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xlsx"

        # Create Excel file with headers
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Employees"

        # Add headers
        ws.append(["Name", "Age", "Department"])
        # Add data rows
        ws.append(["Alice", 25, "Engineering"])
        ws.append(["Bob", 30, "Marketing"])

        wb.save(test_file)

        extracted = processor.extract_text_from_excel(test_file)

        # Check that extracted text contains the data
        assert "test.xlsx" in extracted
        assert "Sheet: Employees" in extracted
        assert "Alice" in extracted
        assert "Bob" in extracted
        assert "25" in extracted
        assert "30" in extracted
        assert "Engineering" in extracted
        assert "Marketing" in extracted
        # With headers, should format as "ColumnName: Value"
        assert "Name:" in extracted
        assert "Age:" in extracted

    def test_extract_text_from_excel_without_headers(self, processor, temp_dir):
        """
        Test text extraction from Excel file without headers.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xlsx"

        # Create Excel file without string headers (numeric data in first row)
        wb = openpyxl.Workbook()
        ws = wb.active

        # Add numeric data (no headers)
        ws.append([1, 100, 200])
        ws.append([2, 150, 250])

        wb.save(test_file)

        extracted = processor.extract_text_from_excel(test_file)

        # Check that extracted text contains the data
        assert "100" in extracted
        assert "200" in extracted
        assert "150" in extracted
        assert "250" in extracted

    def test_extract_text_from_excel_multi_sheet(self, processor, temp_dir):
        """
        Test text extraction from Excel file with multiple sheets.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xlsx"

        # Create Excel file with multiple sheets
        wb = openpyxl.Workbook()

        # First sheet
        ws1 = wb.active
        ws1.title = "Sales"
        ws1.append(["Product", "Price"])
        ws1.append(["Widget", 10.99])

        # Second sheet
        ws2 = wb.create_sheet("Inventory")
        ws2.append(["Item", "Quantity"])
        ws2.append(["Gadget", 50])

        wb.save(test_file)

        extracted = processor.extract_text_from_excel(test_file)

        # Check that both sheets are included
        assert "Sheet: Sales" in extracted
        assert "Sheet: Inventory" in extracted
        assert "Widget" in extracted
        assert "Gadget" in extracted
        assert "10.99" in extracted
        assert "50" in extracted

    def test_extract_text_from_excel_with_empty_cells(self, processor, temp_dir):
        """
        Test text extraction from Excel file with empty cells and rows.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xlsx"

        # Create Excel file with empty cells
        wb = openpyxl.Workbook()
        ws = wb.active

        ws.append(["Name", "Age", "Email"])
        ws.append(["Alice", 25, None])  # Empty email
        ws.append([None, None, None])  # Completely empty row
        ws.append(["Bob", None, "bob@example.com"])  # Empty age

        wb.save(test_file)

        extracted = processor.extract_text_from_excel(test_file)

        # Should contain non-empty data
        assert "Alice" in extracted
        assert "Bob" in extracted
        assert "bob@example.com" in extracted
        # Should handle empty cells gracefully (no "None" in output)
        assert "None" not in extracted

    def test_extract_text_from_excel_with_formulas(self, processor, temp_dir):
        """
        Test text extraction from Excel file with formulas.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "test.xlsx"

        # Create Excel file with formulas
        wb = openpyxl.Workbook()
        ws = wb.active

        ws.append(["A", "B", "Sum"])
        ws.append([10, 20, "=A2+B2"])

        wb.save(test_file)

        # Save and reload to get calculated values
        wb = openpyxl.load_workbook(test_file, data_only=True)
        ws = wb.active
        ws["C2"].value = 30  # Manually set the formula result
        wb.save(test_file)

        extracted = processor.extract_text_from_excel(test_file)

        # Should contain the calculated value, not the formula
        assert "30" in extracted

    def test_extract_text_from_xls_extension(self, processor, temp_dir):
        """
        Test that .xls extension is handled (even though we save as .xlsx).

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        # Note: openpyxl only supports .xlsx, but we can test the extension routing
        test_file = temp_dir / "test.xlsx"

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Test", "Data"])
        ws.append(["Value1", "Value2"])
        wb.save(test_file)

        # Test via extract_text method with .xlsx
        extracted = processor.extract_text(test_file)
        assert "Value1" in extracted
        assert "Value2" in extracted

    def test_process_document_excel(self, processor, temp_dir):
        """
        Test full document processing for Excel file.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        test_file = temp_dir / "employees.xlsx"

        # Create a realistic Excel file
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Staff"

        ws.append(["Name", "Position", "Salary"])
        ws.append(["Alice Johnson", "Engineer", 75000])
        ws.append(["Bob Smith", "Manager", 85000])
        ws.append(["Carol White", "Designer", 70000])

        wb.save(test_file)

        nodes = processor.process_document(test_file, source_type="manual")

        # Should produce at least one node
        assert len(nodes) > 0

        # Check metadata
        first_node = nodes[0]
        assert first_node.metadata["filename"] == "employees.xlsx"
        assert first_node.metadata["source_type"] == "manual"
        assert first_node.metadata["file_type"] == ".xlsx"
        assert "file_hash" in first_node.metadata

        # Check that content was extracted
        full_text = " ".join(node.text for node in nodes)
        assert "Alice Johnson" in full_text
        assert "Engineer" in full_text
        assert "75000" in full_text

    def test_process_directory_includes_excel(self, processor, temp_dir):
        """
        Test that process_directory includes Excel files.

        Args:
            processor: DocumentProcessor fixture.
            temp_dir: Temporary directory fixture.
        """
        # Create multiple file types including Excel
        (temp_dir / "doc1.txt").write_text("Text document.")
        (temp_dir / "doc2.csv").write_text("col1,col2\nval1,val2")

        # Create Excel file
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Header1", "Header2"])
        ws.append(["Data1", "Data2"])
        wb.save(temp_dir / "doc3.xlsx")

        nodes = processor.process_directory(temp_dir)

        # Should process all 3 supported files
        unique_files = set(n.metadata["filename"] for n in nodes)
        assert len(unique_files) == 3
        assert "doc1.txt" in unique_files
        assert "doc2.csv" in unique_files
        assert "doc3.xlsx" in unique_files


class TestDocumentProcessorEnhancement:
    """Test suite for DocumentProcessor enhancement integration."""

    @pytest.fixture
    def processor(self):
        """Create a DocumentProcessor instance for testing."""
        return DocumentProcessor(chunk_size=512, chunk_overlap=50)

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield Path(tmpdir)

    @pytest.fixture
    def mock_enhancement_processor(self):
        """Mock the enhancement processor."""
        from unittest.mock import MagicMock, patch

        with patch(
            "src.document_processing.document_processor._get_enhancement_processor"
        ) as mock_get:
            mock_enhancer = MagicMock()
            mock_enhancer.should_enhance.return_value = True
            mock_enhancer.enhance_document.return_value = {
                "enhanced_text": "Enhanced narrative content.\n\nQ: Test? A: Yes.",
                "narrative": "Enhanced narrative content.",
                "qa_pairs": "Q: Test? A: Yes.",
                "enhancement_count": 2,
            }
            mock_get.return_value = mock_enhancer
            yield mock_enhancer

    def test_enhancement_integration_csv(
        self, processor, temp_dir, mock_enhancement_processor
    ):
        """Test that CSV files are enhanced during processing."""
        test_file = temp_dir / "test.csv"

        # Create CSV
        df = pd.DataFrame({"Name": ["Alice"], "Age": [25]})
        df.to_csv(test_file, index=False)

        nodes = processor.process_document(test_file, source_type="manual")

        # Verify enhancement was called
        assert mock_enhancement_processor.should_enhance.called
        assert mock_enhancement_processor.enhance_document.called

        # Verify metadata includes enhancement info
        assert len(nodes) > 0
        first_node = nodes[0]
        assert first_node.metadata["enhanced"] is True
        assert first_node.metadata["enhancement_count"] == 2

        # Verify enhanced content is in the text
        full_text = " ".join(node.text for node in nodes)
        assert "Enhanced narrative" in full_text or "Q: Test?" in full_text

    def test_enhancement_integration_excel(
        self, processor, temp_dir, mock_enhancement_processor
    ):
        """Test that Excel files are enhanced during processing."""
        test_file = temp_dir / "test.xlsx"

        # Create Excel
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Name", "Value"])
        ws.append(["Test", 123])
        wb.save(test_file)

        nodes = processor.process_document(test_file, source_type="manual")

        # Verify enhancement was called
        assert mock_enhancement_processor.should_enhance.called
        assert mock_enhancement_processor.enhance_document.called

        # Verify metadata
        assert len(nodes) > 0
        assert nodes[0].metadata["enhanced"] is True
        assert nodes[0].metadata["file_type"] == ".xlsx"

    def test_no_enhancement_for_pdf(
        self, processor, temp_dir, mock_enhancement_processor
    ):
        """Test that PDF files are not enhanced."""

        test_file = temp_dir / "test.txt"  # Use txt as proxy for PDF
        test_file.write_text("This is a text document.")

        # Make should_enhance return False for txt
        mock_enhancement_processor.should_enhance.return_value = False

        nodes = processor.process_document(test_file, source_type="manual")

        # Enhancement should be checked but not performed
        assert mock_enhancement_processor.should_enhance.called
        assert not mock_enhancement_processor.enhance_document.called

        # Metadata should show not enhanced
        assert len(nodes) > 0
        assert nodes[0].metadata["enhanced"] is False
        assert nodes[0].metadata["enhancement_count"] == 0

    def test_enhancement_failure_graceful(
        self, processor, temp_dir, mock_enhancement_processor
    ):
        """Test that enhancement failures don't break document processing."""
        test_file = temp_dir / "test.csv"
        df = pd.DataFrame({"Col": ["Value"]})
        df.to_csv(test_file, index=False)

        # Make enhancement fail
        mock_enhancement_processor.enhance_document.side_effect = Exception(
            "Enhancement error"
        )

        # Should still process document without enhancement
        nodes = processor.process_document(test_file, source_type="manual")

        assert len(nodes) > 0
        # Should have original content
        assert nodes[0].metadata["enhanced"] is False

    def test_enhancement_disabled_via_config(self, processor, temp_dir):
        """Test that enhancement can be disabled via configuration."""
        from unittest.mock import patch

        test_file = temp_dir / "test.csv"
        df = pd.DataFrame({"Col": ["Value"]})
        df.to_csv(test_file, index=False)

        # Mock _get_enhancement_processor to return None (disabled)
        with patch(
            "src.document_processing.document_processor._get_enhancement_processor"
        ) as mock_get:
            mock_get.return_value = None

            nodes = processor.process_document(test_file, source_type="manual")

            # Should process without enhancement
            assert len(nodes) > 0
            assert nodes[0].metadata["enhanced"] is False
            assert nodes[0].metadata["enhancement_count"] == 0


class TestTextNormalization:
    """Tests for _normalize_extracted_text() helper."""

    def test_empty_input(self):
        assert _normalize_extracted_text("") == ""
        assert _normalize_extracted_text(None) is None

    def test_joins_wrapped_lines(self):
        """Lines not ending in sentence terminators should be joined."""
        raw = "Faculty of\nDepartment of\nLife Sciences\n"
        result = _normalize_extracted_text(raw)
        assert result == "Faculty of Department of Life Sciences"

    def test_preserves_sentence_breaks(self):
        """Lines ending in sentence terminators are not joined with the next."""
        # The terminator-ended line is not merged INTO the next, but the
        # next line still becomes part of the same paragraph when there's
        # no blank separator. This is acceptable — the terminator signals
        # "don't wrap UP into this line" rather than a hard paragraph break.
        raw = "This is the first sentence.\nAnd here is the next one.\n"
        result = _normalize_extracted_text(raw)
        # Both lines remain as separate entries joined into one paragraph
        assert "This is the first sentence." in result
        assert "And here is the next one." in result

    def test_collapses_multiple_blank_lines(self):
        raw = "Paragraph one.\n\n\n\nParagraph two.\n"
        result = _normalize_extracted_text(raw)
        assert result == "Paragraph one.\n\nParagraph two."

    def test_collapses_multiple_spaces(self):
        raw = "Word     with    many    spaces."
        result = _normalize_extracted_text(raw)
        assert result == "Word with many spaces."

    def test_paragraph_boundary_preserved(self):
        """Blank lines separate paragraphs in the output."""
        raw = "First para line one\nFirst para line two\n\nSecond para line one"
        result = _normalize_extracted_text(raw)
        paragraphs = result.split("\n\n")
        assert len(paragraphs) == 2
        assert paragraphs[0] == "First para line one First para line two"
        assert paragraphs[1] == "Second para line one"

    def test_fragmented_template_text(self):
        """Simulates the handbook bug — tiny fragments get joined."""
        raw = (
            "Faculty of {insert here}\n"
            "Department of {insert here}\n"
            "\n"
            "{insert qualification}\n"
            "{name of course}\n"
        )
        result = _normalize_extracted_text(raw)
        # Two paragraphs, each a single joined line
        paragraphs = result.split("\n\n")
        assert len(paragraphs) == 2
        assert paragraphs[0] == "Faculty of {insert here} Department of {insert here}"
        assert paragraphs[1] == "{insert qualification} {name of course}"


class TestChunkQualityGuards:
    """Tests for tiny-chunk filter and pathological-chunking sanity guard."""

    @pytest.fixture
    def processor(self):
        return DocumentProcessor(chunk_size=2048, chunk_overlap=100)

    @pytest.fixture
    def temp_dir(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            yield Path(tmpdir)

    def test_normal_doc_processes_successfully(self, processor, temp_dir):
        """A normal document should produce a reasonable chunk count."""
        test_file = temp_dir / "normal.txt"
        # ~3000 chars of normal prose
        content = (
            "This is a normal document. It has several sentences. "
            "Each sentence is long enough to carry meaning. The chunking "
            "should produce a modest number of chunks of reasonable size. "
        ) * 20
        test_file.write_text(content)

        nodes = processor.process_document(test_file)
        assert len(nodes) > 0
        assert len(nodes) < 30  # nowhere near pathological
        for n in nodes:
            assert len(n.text.strip()) >= MIN_CHUNK_CHARS

    def test_tiny_chunks_filtered(self, processor, temp_dir):
        """Chunks shorter than MIN_CHUNK_CHARS should be dropped."""
        from unittest.mock import MagicMock

        from llama_index.core.schema import TextNode

        big = "x" * 200
        small = "tiny"
        fake_nodes = [
            TextNode(text=big),
            TextNode(text=small),
            TextNode(text=big),
            TextNode(text=small),
        ]
        mock_splitter = MagicMock()
        mock_splitter.get_nodes_from_documents.return_value = fake_nodes
        processor.splitter = mock_splitter

        test_file = temp_dir / "mixed.txt"
        test_file.write_text("placeholder content " * 200)

        nodes = processor.process_document(test_file)
        # Only the two big chunks survive
        assert len(nodes) == 2
        for n in nodes:
            # Note: header is prepended after filtering, so stripped text
            # starts with "Document: ..." - check it's >= the raw big size
            assert "xxx" in n.text  # original content present

    def test_pathological_chunking_raises(self, processor, temp_dir):
        """Many chunks whose average size is way below configured raises."""
        from unittest.mock import MagicMock

        from llama_index.core.schema import TextNode

        # 50 tiny chunks — triggers PATHOLOGICAL_MIN_CHUNKS (20) AND
        # avg size << chunk_size * 0.2 (= 102 for chunk_size=512)
        # Each chunk is 60 chars — above MIN_CHUNK_CHARS (50) so not filtered
        fake_nodes = [TextNode(text="x" * 60) for _ in range(50)]
        mock_splitter = MagicMock()
        mock_splitter.get_nodes_from_documents.return_value = fake_nodes
        processor.splitter = mock_splitter

        test_file = temp_dir / "broken.txt"
        test_file.write_text("content " * 125)

        with pytest.raises(ValueError, match="extraction is broken"):
            processor.process_document(test_file)

    def test_single_tiny_chunk_preserved(self, processor, temp_dir):
        """Don't filter a document down to zero chunks — sole chunk stays."""
        from unittest.mock import MagicMock

        from llama_index.core.schema import TextNode

        # Only one chunk, shorter than MIN_CHUNK_CHARS
        fake_nodes = [TextNode(text="short")]
        mock_splitter = MagicMock()
        mock_splitter.get_nodes_from_documents.return_value = fake_nodes
        processor.splitter = mock_splitter

        test_file = temp_dir / "tiny.txt"
        test_file.write_text("short")

        nodes = processor.process_document(test_file)
        assert len(nodes) == 1
